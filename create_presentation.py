"""
PowerPoint Presentation Generator for Inventory and Sales Management System
Usage: python create_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

print("Creating PowerPoint Presentation...")

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Colors
PRIMARY = RGBColor(79, 70, 229)
SECONDARY = RGBColor(22, 163, 74)
ACCENT = RGBColor(245, 158, 11)
TEXT = RGBColor(31, 41, 55)
WHITE = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(248, 250, 252)

print("Adding Slide 1: Title Slide...")
# SLIDE 1: Title Slide
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide1.background
bg.fill.solid()
bg.fill.fore_color.rgb = PRIMARY

title = slide1.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(1))
tf = title.text_frame
tf.text = "Inventory and Sales Management System"
p = tf.paragraphs[0]
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

subtitle = slide1.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.6))
tf = subtitle.text_frame
tf.text = "Diploma in Software Engineering | Semester 4"
p = tf.paragraphs[0]
p.font.size = Pt(22)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

team = slide1.shapes.add_textbox(Inches(2), Inches(5), Inches(6), Inches(1.5))
tf = team.text_frame
tf.text = "GROUP MEMBERS:\n• Nakibinge Collins - Backend Operations\n• Sekimpi Ibrahim - Frontend Operations\n• Lwebugga Aaron - Database Design"
for para in tf.paragraphs:
    para.font.size = Pt(16)
    para.font.color.rgb = WHITE
    para.alignment = PP_ALIGN.CENTER

print("Adding Slide 2: Problem Statement...")
# SLIDE 2: Problem Statement
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

title = slide2.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
tf = title.text_frame
tf.text = "Problem Statement"
p = tf.paragraphs[0]
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = PRIMARY

# Left side
left = slide2.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(4.3), Inches(5.5))
tf = left.text_frame
tf.text = "What Problems Does This Software Solve?\n\nReal-World Challenges:\n📝 Manual inventory tracking\n⚠️ Stock running out without warning\n💔 Lost sales opportunities\n📊 No financial visibility\n👥 No role separation\n🏭 Poor supplier management"
for para in tf.paragraphs:
    para.font.size = Pt(14) if para.text.startswith("What") else Pt(16)
    para.font.color.rgb = TEXT
    para.space_before = Pt(8)
    if para.text.startswith("What"):
        para.font.bold = True
        para.font.size = Pt(18)

# Right side
right = slide2.shapes.add_textbox(Inches(5.1), Inches(1.3), Inches(4.3), Inches(5.5))
tf = right.text_frame
tf.text = "Who Needs This Software?\n\n✓ Small retail shops\n✓ Pharmacies\n✓ Hardware stores\n✓ Bookshops\n✓ Wholesale businesses\n✓ Supermarkets\n\nTarget:\n1-50 employees"
for para in tf.paragraphs:
    para.font.size = Pt(16)
    para.font.color.rgb = TEXT
    para.space_before = Pt(8)
    if para.text.startswith("Who"):
        para.font.bold = True
        para.font.size = Pt(18)


print("Adding Slide 3: System Overview...")
# SLIDE 3: System Overview
slide3 = prs.slides.add_slide(prs.slide_layouts[6])

title = slide3.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
tf = title.text_frame
tf.text = "System Overview & Key Features"
p = tf.paragraphs[0]
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = PRIMARY

desc = slide3.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(8.8), Inches(0.6))
tf = desc.text_frame
tf.text = "A modern system that automates business operations from stock tracking to sales processing and financial reporting."
p = tf.paragraphs[0]
p.font.size = Pt(14)
p.font.color.rgb = TEXT

# Feature boxes - 2x3 grid
features = [
    ("📦 Inventory", "• Real-time tracking\n• Low-stock alerts\n• Categorization"),
    ("💰 Sales", "• POS system\n• Multiple payments\n• Receipt generation"),
    ("🛒 Purchases", "• Record purchases\n• Auto stock updates\n• Supplier tracking"),
    ("📊 Reports", "• Daily sales\n• Stock movements\n• AI-powered insights"),
    ("👥 Users", "• Multi-tenant\n• Role-based access\n• Secure auth"),
    ("🤖 AI Assistant", "• Mistral AI\n• Natural language\n• Predictions")
]

x_positions = [0.6, 3.6, 6.6]
y_positions = [2.2, 4.4]
idx = 0

for row in y_positions:
    for col in x_positions:
        box = slide3.shapes.add_shape(1, Inches(col), Inches(row), Inches(2.8), Inches(2))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BG
        box.line.color.rgb = PRIMARY
        
        tf = box.text_frame
        tf.text = f"{features[idx][0]}\n\n{features[idx][1]}"
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = PRIMARY
        
        for para in tf.paragraphs[1:]:
            para.font.size = Pt(12)
            para.font.color.rgb = TEXT
        
        idx += 1
        if idx >= len(features):
            break

print("Adding Slide 4: Architecture & Design...")
# SLIDE 4: Architecture & Design  
slide4 = prs.slides.add_slide(prs.slide_layouts[6])

title = slide4.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
tf = title.text_frame
tf.text = "Architecture & Design"
p = tf.paragraphs[0]
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = PRIMARY

# Architecture diagram (left)
arch = slide4.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(4.3), Inches(5.5))
tf = arch.text_frame
tf.text = "SYSTEM ARCHITECTURE\n\n┌─────────────┐\n│ React Frontend │\n│  Port 3000   │\n└──────┬──────┘\n       ↓\n   REST API\n       ↓\n┌──────┴──────┐\n│Laravel Backend│\n│  35+ Endpoints│\n└──────┬──────┘\n       ↓\n  Eloquent ORM\n       ↓\n┌──────┴──────┐\n│MySQL Database│\n│  14 Tables   │\n└─────────────┘"
for para in tf.paragraphs:
    para.font.size = Pt(11)
    para.font.name = "Courier New"
    para.font.color.rgb = TEXT
    if para.text.startswith("SYSTEM"):
        para.font.bold = True
        para.font.size = Pt(14)
        para.font.color.rgb = PRIMARY

# Technologies (right)
tech = slide4.shapes.add_textbox(Inches(5.1), Inches(1.3), Inches(4.3), Inches(5.5))
tf = tech.text_frame
tf.text = "TECHNOLOGIES USED\n\nBackend:\n• PHP 8.2+\n• Laravel 12\n• MySQL 8.0+\n• Laravel Sanctum\n\nFrontend:\n• React 19\n• JavaScript ES6+\n• HTML5 & CSS3\n\nTools:\n• Git & GitHub\n• VS Code\n• Postman\n• XAMPP/WAMP\n\nAI:\n• Mistral AI"
for para in tf.paragraphs:
    para.font.size = Pt(14)
    para.font.color.rgb = TEXT
    if para.text in ["TECHNOLOGIES USED", "Backend:", "Frontend:", "Tools:", "AI:"]:
        para.font.bold = True
        para.font.color.rgb = PRIMARY if para.text == "TECHNOLOGIES USED" else SECONDARY

print("Adding Slide 5: Challenges & Solutions...")
# SLIDE 5: Challenges & Solutions
slide5 = prs.slides.add_slide(prs.slide_layouts[6])

title = slide5.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
tf = title.text_frame
tf.text = "Challenges & Solutions"
p = tf.paragraphs[0]
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = PRIMARY

content = slide5.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(8.8), Inches(5.5))
tf = content.text_frame
tf.text = """Challenge 1: Different Business Requirements
• Surveyed 20 businesses across different sectors
• Identified common pain points (80% coverage)
• Designed flexible, extensible architecture

Challenge 2: Multi-Tenant Data Isolation
• Implemented tenant_id in all tables
• Automatic tenant scoping & middleware
• Database-level constraints

Challenge 3: Real-Time Stock Management
• Database transactions for atomic operations
• Stock validation before sales
• Complete audit trail

Challenge 4: User Role Complexity
• Role-based access control (RBAC)
• 20+ granular permissions
• Frontend & backend authorization

Survey Results:
✓ 85% track inventory manually
✓ 70% experienced stockouts
✓ 60% cannot generate reports easily
✓ 90% want mobile access"""

for para in tf.paragraphs:
    para.font.size = Pt(13)
    para.font.color.rgb = TEXT
    para.space_before = Pt(6)
    if para.text.startswith("Challenge"):
        para.font.bold = True
        para.font.size = Pt(15)
        para.font.color.rgb = PRIMARY
    elif para.text.startswith("Survey"):
        para.font.bold = True
        para.font.size = Pt(15)
        para.font.color.rgb = SECONDARY


print("Adding Slide 6: Conclusion...")
# SLIDE 6: Conclusion
slide6 = prs.slides.add_slide(prs.slide_layouts[6])

title = slide6.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
tf = title.text_frame
tf.text = "Conclusion & Future Vision"
p = tf.paragraphs[0]
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = PRIMARY

# Left - Achievements
left = slide6.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(4.3), Inches(5.5))
tf = left.text_frame
tf.text = """KEY ACHIEVEMENTS

✓ Complete Backend API
✓ Modern Frontend
✓ Multi-Tenant System
✓ Role-Based Access
✓ Real-Time Stock Mgmt
✓ AI-Powered Analytics
✓ Comprehensive Reports

IMPACT:
• 70% less manual work
• Zero stock errors
• Real-time visibility
• Data-driven decisions"""

for para in tf.paragraphs:
    para.font.size = Pt(14)
    para.font.color.rgb = TEXT
    if para.text in ["KEY ACHIEVEMENTS", "IMPACT:"]:
        para.font.bold = True
        para.font.size = Pt(16)
        para.font.color.rgb = PRIMARY

# Right - Future
right = slide6.shapes.add_textbox(Inches(5.1), Inches(1.3), Inches(4.3), Inches(5.5))
tf = right.text_frame
tf.text = """FUTURE IMPROVEMENTS

Phase 1 (3-6 months):
• Barcode scanning
• Email notifications
• Advanced reports
• Mobile app

Phase 2 (6-12 months):
• Loyalty program
• Multi-currency
• E-commerce
• Auto-reordering

Phase 3 (1-2 years):
• Business intelligence
• Multi-location
• Accounting integration
• Public API"""

for para in tf.paragraphs:
    para.font.size = Pt(13)
    para.font.color.rgb = TEXT
    if para.text in ["FUTURE IMPROVEMENTS", "Phase 1 (3-6 months):", "Phase 2 (6-12 months):", "Phase 3 (1-2 years):"]:
        para.font.bold = True
        para.font.color.rgb = SECONDARY if not para.text.startswith("FUTURE") else PRIMARY

print("Adding Slide 7: Thank You...")
# SLIDE 7: Thank You
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide7.background
bg.fill.solid()
bg.fill.fore_color.rgb = SECONDARY

thank = slide7.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(1.5))
tf = thank.text_frame
tf.text = "THANK YOU!"
p = tf.paragraphs[0]
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

sub = slide7.shapes.add_textbox(Inches(2), Inches(4.2), Inches(6), Inches(1))
tf = sub.text_frame
tf.text = "Questions & Demonstration"
p = tf.paragraphs[0]
p.font.size = Pt(28)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Save presentation
filename = "INVENTORY_SALES_PRESENTATION.pptx"
prs.save(filename)
print(f"\n✓ PowerPoint presentation created successfully!")
print(f"✓ File saved as: {filename}")
print(f"\nOpen the file in Microsoft PowerPoint to view and edit.")
